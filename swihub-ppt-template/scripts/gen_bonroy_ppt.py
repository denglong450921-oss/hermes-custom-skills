#!/usr/bin/env python3
"""Convert 员工手册制作.md → Bonroy Apex Sync presentation using swihub template."""

import copy, os, re
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

TEMPLATE = "/Users/f/.hermes/skills/swihub-ppt-template/assets/template.pptx"
OUTPUT   = "/Users/f/Documents/github项目2026/ppt_Bonroy_Apex_Sync/projects/bonroy_apex_sync_ppt169_20260519/员工手册制作.pptx"

prs = Presentation(TEMPLATE)

# ── helpers ──
def find_shape(slide, name):
    for s in slide.shapes:
        if s.name == name:
            return s
    return None

def find_shape_by_y(slide, tag_y):
    """Find shape with tag like 'Text 2@1.10' = name Text 2, min_y ~1.10."""
    name, y_str = tag_y.rsplit('@', 1)
    target_y = float(y_str)
    best = None; best_d = 999
    for s in slide.shapes:
        if s.name == name:
            d = abs(s.top/914400 - target_y)
            if d < best_d:
                best_d = d; best = s
    return best

def replace_text_on_slide(slide, mapping):
    """For each shape on slide, if its text contains a key in mapping, replace whole text."""
    for s in slide.shapes:
        if not s.has_text_frame: continue
        t = s.text_frame.text.strip()
        if not t: continue  # skip empty shapes (background cards, decorative bars)
        for orig, new in mapping.items():
            if t == orig or t.startswith(orig[:40]) or (len(orig) > 30 and orig.startswith(t[:40])):
                s.text = new
                break

def fill_table(slide, table_name, row_data_list):
    """Fill table cells with row_data_list: [(row, col, text), ...]"""
    for shape in slide.shapes:
        if shape.has_table and shape.name == table_name:
            tbl = shape.table
            for r, c, txt in row_data_list:
                if r < len(tbl.rows) and c < len(tbl.columns):
                    tbl.cell(r, c).text = txt
            return

def duplicate_slide(prs, slide_index):
    ts = prs.slides[slide_index]
    ns = prs.slides.add_slide(ts.slide_layout)
    for shape in list(ns.shapes):
        sp = shape._element; sp.getparent().remove(sp)
    rmap = {}
    for rel in ts.part.rels.values():
        if rel.is_external:
            rmap[rel.rId] = ns.part.relate_to(rel.target_ref, rel.reltype)
        else:
            rmap[rel.rId] = ns.part.relate_to(rel.target_part, rel.reltype)
    tspt = ts._element.find(qn('p:cSld')).find(qn('p:spTree'))
    nspt = ns._element.find(qn('p:cSld')).find(qn('p:spTree'))
    for child in list(tspt):
        nc = copy.deepcopy(child)
        for attr in ['r:embed','r:link','r:id']:
            qa = qn(attr)
            for e in nc.iter():
                if qa in e.attrib and e.attrib[qa] in rmap:
                    e.attrib[qa] = rmap[e.attrib[qa]]
        nspt.append(nc)
    return ns

# ── Content Mapping ──
# Each entry: (template_idx, replace_dict)
# Keys are EXACT original text or FIRST 40+ chars (enough to match uniquely)

slides_map = [
# ═══ Slide 0: Cover ═══
(0, {
    "Swihub Solutions LLC.": "Bonroy Apex Sync",
    "Actualización Estratégica": "战略升级",
    "Remodelar el panorama del tráfico global aprovechando el capital de NY.": "利用纽约的资本与资源优势，重塑全球电商流量格局",
}),

# ═══ Slide 1: Mission & Vision ═══
(1, {
    "Misión y Visión": "使命与愿景",
    "Misión Principal": "核心使命",
    "Empoderar a los usuarios globales para lograr un crecimiento de ingresos digitales": "赋能全球用户，实现数字收入增长",
    "Nuestra Creencia": "我们的信念",
    "La tecnología no debería servir solo a las grandes empresas. Nos comprometemos a permitir que los usuarios comunes parti": "技术不应只服务于大型企业。我们致力于让全球普通用户，公平参与数字经济并获得收益",
    "Capacidades Centrales": "核心能力",
    "Influencia Global:": "全球影响力",
    "Cubriendo mercados globales, logrando un crecimiento en múltiples regiones": "覆盖全球市场，实现跨时区、跨区域增长",
    "Automatización Inteligente:": "智能自动化",
    "Automatización de procesos y mejora de eficiencia basadas en IA": "基于AI实现业务流程自动化与效率提升",
    "Centrado en el Usuario:": "用户为中心",
    "Enfatizando la participación, haciendo de cada usuario un creador de valor": "强调用户参与，让每个用户都成为价值创造者",
    "Que cada persona común\npueda obtener crecimiento de ingresos\nen el comercio digital global": "让每一个普通人，都能在全球数字商业中获得收入增长",
}),

# ═══ Slide 2: CSR ═══
(2, {
    "Responsabilidad Social Corporativa": "企业社会责任",
    "Concepto Central: Como empresa responsable, nos comprometemos a retribuir a la sociedad.": "核心理念：以责任为灯，照亮他人，也照亮自己前行的路",
    "Nuestro Compromiso\nSwihub no solo persigue el éxito comercial, sino también la creación de valor social a largo plazo:\nA": "我们的承诺\nBonroy Apex Sync在追求商业卓越的同时，始终着眼于长期社会价值的创造",
    "Direcciones de Acción Clave": "核心行动方向",
    "Donación Pública": "公益捐赠",
    "Proporcionar apoyo financiero para proyectos de bienestar público": "提供资金支持，推动社会公益项目发展",
    "Empoderamiento Edu.": "教育赋能",
    "Promover programas de intercambio de conocimientos y formación": "推动知识共享与人才培养计划",
    "Resp. Ambiental": "环境责任",
    "Apoyar el desarrollo sostenible y acciones ambientales": "支持可持续发展与环保行动",
    "Construcción Comunitaria": "社区建设",
    "Fortalecer la conexión y participación de la comunidad local": "加强本地社区连接与参与",
    "Áreas de Inversión Continua\nActividades de Bienestar: Realizar acciones de bienestar social regularmente\nImpacto Social:": "持续投入领域\n公益活动：定期开展社会公益行动\n社会影响：促进社区成长与社会进步",
    "Conclusión: Equilibrando el crecimiento comercial y el valor social, construyendo un impacto sostenible.": "商业增长与社会价值并重，构建长期可持续影响力",
}),

# ═══ Slide 3: Philanthropy ═══
(3, {
    "Espíritu Filantrópico": "慈善精神",
    "Concepto Central: La filantropía no es solo donar, sino transmitir valores": "核心理念：我们把心意变成捐赠，更把信念变成传递",
    "Nuestra Comprensión": "我们的理解",
    "La verdadera filantropía no es solo transferir recursos, sino crear un impacto a largo plazo:\n\nConectando a donantes y b": "真正的慈善，不只是资源的转移，而是创造长期影响力：\n\n连接施予者与受益者\n在社会中形成正向循环\n放大个体与群体的社会价值",
    "Valores Fundamentales": "核心价值观",
    "Compasión:": "同情",
    "Entender situaciones ajenas": "理解他人处境",
    "Responsabilidad:": "责任",
    "Asumir obligaciones sociales": "承担社会义务",
    "Empatía:": "共情",
    "Acción y conexión proactivas": "主动行动与连接",
    "Nuestro Camino Práctico": "我们的实践路径",
    "Impulsando un cambio social a largo plazo combinando caridad y educación:\n\nCultivando personas con empatía y responsabil": "通过将慈善与教育结合，推动长期社会改变：\n\n培养具备同理心与责任感的个体\n引导年轻一代形成正确的价值观\n鼓励主动参与社会与公共事务",
    "Áreas de Impacto Clave": "重点影响领域",
    "Edu. y Crecimiento Juvenil": "青少年教育与成长",
    "Desarrollo Comunitario": "社区发展",
    "Alfabetización Digital": "数字素养与认知提升",
    "Oportunidades Económicas": "经济机会与赋能",
    "Conclusión: Impulsar filantropía con valores, amplificando retornos sociales con impacto.": "用价值观驱动慈善，用长期影响力放大社会回报",
}),

# ═══ Slide 4: Milestone ═══
(4, {
    "Hito: Adquisición de Capital K Ltd.": "里程碑：收购 Capital K Ltd.",
    "Evento Clave\nCompletó la adquisición de Capital K Ltd., una firma de inversión estadounidense, marcando la integración o": "关键事件\n完成对美国成熟投资机构 Capital K Ltd.的收购\n标志着平台正式引入机构级金融能力",
    "Importancia Estratégica y Mejora de Capacidad": "战略意义与能力升级",
    "Respaldo Crediticio Institucional:\nIntroducir sistema financiero maduro para mejorar la confianza (Fondos asegurados)\n\nC": "机构信用背书：引入成熟金融体系与专业经验\n风险控制能力：建立更完善的风控与资产管理体系\n合规能力提升：强化监管标准",
    "Mejora de capacidades (Valor práctico)": "能力升级（落地价值）",
    "Mayor seguridad\nSeguridad institucional para el fondo común\n\n\nLiquidación más estable\nFondos de usuario protegidos por u": "更高安全性：资金池获得机构级保障\n更稳健结算：用户资金由专业风控体系保护\n专业化管理：引入长期资产管理经验",
    "Conclusión: De \"plataforma\" a \"infraestructura global de activos con capacidades financieras institucionales\".": "从「平台」升级为「具备机构级金融能力的全球资产与流量基础设施」",
}),

# ═══ Slide 5: HQ ═══
(5, {
    "Sede Central: Nueva York": "公司总部：伦敦",
    "Ubicación Estratégica: Midtown Manhattan, Nueva York, EE. UU. — en el núcleo financiero, conectando capital mundial.": "战略选址：英国伦敦——位于全球金融核心枢纽，直接连接世界资本与资源",
    "Importancia Estratégica\nEstablecer la sede en Midtown Manhattan para construir el centro de control global:\nCerca del nú": "战略意义\n将总部设立于伦敦，构建全球业务的核心控制中枢：\n\n靠近全球金融体系核心\n强化国际化运营能力\n提升金融合规与监管对接能力",
    "Conclusión: \"En el centro financiero global, construyendo una infraestructura de comercio digital mundial.\"": "「站在全球金融中心，构建面向世界的数字商业基础设施」",
}),

# ═══ Slide 6: London Ops → Europe Ops Center ═══
(6, {
    "Centro de Ops América N/S: Londres": "欧洲运营中心：伦敦",
    "Ubicación: Londres, Reino Unido": "地点：英国伦敦",
    "Descripción: Ubicados en Londres, ciudad central de Europa, conectamos mercados globales.": "我们位于欧洲核心城市伦敦，连接全球市场，提升业务效率。",
    "Ventajas Principales:": "核心优势：",
    "Puerta Europea": "欧洲门户",
    "Operaciones Cumplimiento": "合规运营",
    "Eficiencia Colaborativa": "协同高效",
    "Apoyo Local": "本地支持",
}),

# ═══ Slide 7: Pain Points ═══
(7, {
    "Puntos Débiles: Desafíos del E-commerce Tradicional": "当前行业痛点：传统电子商务面临的挑战",
    "Problemas Centrales: Altos costos, umbrales y riesgos limitan la participación.": "核心问题：高成本、高门槛、高风险，限制了大规模参与与增长",
    "Tres Desafíos Principales:": "三大核心挑战",
    "Costos de Tráfico": "高流量成本",
    "La adquisición de clientes depende de gran inversión publicitaria; ROI decreciente.": "获客依赖大量广告投入\n投入高、回报不确定，ROI持续下降",
    "Altos Umbrales de Entrada": "高进入门槛",
    "La tecnología y las reglas de plataforma dificultan la entrada de usuarios comunes y pymes.": "技术、平台规则与资金要求复杂\n普通用户与中小参与者难以进入",
    "Altos Riesgos Operativos": "高运营风险",
    "Presión de inventario + complejidad logística, gran ocupación de capital.": "库存压力 + 物流复杂性\n资金占用大，风险不可控",
    "Demanda del Mercado:": "市场需求",
    "Los empresarios enfrentan presión de inventario y gestión logística compleja. Se necesita un modelo más eficiente.": "企业家面临着库存压力和复杂的物流管理。市场迫切需要一种更轻便、更高效的参与模式。",
    "Solución Necesaria: Surge un modelo nuevo": "需要解决方案：一个全新的模式正在出现",
    "Participación de bajo umbral": "低门槛参与",
    "Tráfico descentralizado": "去中心化流量",
    "+ Operaciones inteligentes": "+ 智能化运营",
}),

# ═══ Slide 8: Solution ═══
(8, {
    "Solución: Modelo de Alquiler de Tienda Digital": "解决方案：数字店铺租赁模型",
    "Cambiando de \"E-commerce Pesado\" a \"Modelo de Ingresos Ligeros\"": "从「重运营电商」转向「轻参与收益模型」",
    "Cómo Resolvemos los Problemas:\nPara los 3 puntos débiles del e-commerce tradicional, Swihub ofrece soluciones estructura": "我们如何解决问题\n针对传统电商三大痛点，Bonroy Apex Sync提供结构性解决方案：\n\n将复杂运营标准化\n将高门槛参与产品化\n将不确定收益模型数据化",
    "Diseño Central del Producto:": "核心产品设计",
    "Operaciones Estándar": "标准化运营",
    "Transformando procesos complejos en activos digitales replicables.": "将复杂流程转化为可复制的数字资产\n降低运营难度",
    "Tiendas Maduras": "成熟店铺供给",
    "Proveyendo tiendas con tráfico y bienes, los usuarios no empiezan desde cero.": "提供已具备流量与商品的可运营店铺\n用户无需从0开始",
    "Reparto de Ingresos": "收益分成机制",
    "Reparto de beneficios basado en rendimiento real, retornos transparentes.": "基于真实表现进行分润\n收益透明、可量化",
    "Ventajas del Modelo:\nParticipación ligera (Sin inventario ni logística)\nReplicación escalable (Operaciones unificadas)\nI": "商业模式优势\n轻资产参与（用户无需库存与物流）\n可规模化复制（平台统一运营）\n收益与数据挂钩（增强信任与留存）",
    "\"Hemos transformado el comercio electrónico de un 'negocio' a un tipo de 'activo participativo'.\"": "「我们把电商，从一门生意，变成一类可参与的资产」",
}),

# ═══ Slide 9: Guarantee ═══
(9, {
    "Mecanismo de Crédito: Sistema de Depósito": "保障机制：可靠的权益体系",
    "Protección del Consumidor": "消费者保护",
    "Los depósitos de tienda aseguran la protección del consumidor y mejoran la credibilidad.": "我们在坚实的合同基础上运作，确保每一步的透明与互信。",
    "Totalmente Reembolsable": "全额可退",
    "Si decide dejar la plataforma, su depósito se reembolsa totalmente en cualquier momento sin cargos. Sin bloqueo.": "如果您决定离开平台，您可以在缴纳租金后一个月内要求全额退款",
    "Beneficios del Depósito:": "提供租金的收益",
    "💡 Impulso de Tráfico": "流量支持提升",
    "Obtener prioridad en visualización": "在平台中获得优先展示机会",
    "💡 Privilegios Operativos": "运营权限",
    "Desbloquear funciones mejoradas": "解锁增强的平台功能",
    "💡 Mejora de Crédito": "信用提升",
    "Obtener mejores calificaciones de confianza": "获得更高的信任评级",
}),

# ═══ Slide 10: Store Level Overview ═══
(10, {
    "Mecanismo de Crédito: Sistema de Depósito": "店铺等级概览",
    "🚀 Cómo resolvemos el problema:": "三级店铺体系",
    "Swihub ofrece soluciones estructurales a los tres principales problemas del comercio electrónico tradicional:\nEstandariz": "B1 基础店铺 → B2 进阶店铺 → I1 高级店铺\n\n不同等级对应不同租金与收益水平，满足各类用户需求",
    "📈 Ventajas del modelo de negocio": "参与模式优势",
    "Participación mínima en activos (los usuarios no necesitan inventario ni logística) \nReplicación escalable (operación de": "低门槛参与 | 可规模化复制 | 收益透明可量化",
    "Diseño de producto principal": "核心设计理念",
    "💡 Operaciones estandarizadas": "标准化运营",
    "Transformación de procesos complejos en activos digitales replicables Reducción de la complejidad operativa": "将复杂电商流程转化为可复制的数字资产，大幅降低运营难度",
    "💡 Suministro de tiendas consolidadas": "成熟店铺供给",
    "Provisión de tráfico y productos a tiendas operativas existentes Los usuarios no necesitan empezar de cero": "提供已具备流量与商品的运营店铺，用户无需从零开始",
    "💡 Mecanismo de reparto de ingresos": "收益分成机制",
    "Reparto de beneficios basado en el rendimiento real Ingresos transparentes y cuantificables.": "基于真实店铺表现进行分润，收益透明且可量化",
    "\"Hemos transformado el comercio electrónico de un 'negocio' a un tipo de 'activo participativo'.\"": "「我们把电商，从一门生意，变成一类可参与的资产」",
}),

# ═══ Slide 11: Micro Store → B1 ═══
(11, {
    "Plan de Tienda Micro": "B1店铺收益计划",
    "Tienda Micro, ideal para principiantes": "轻松起步，最佳选择",
    "🔐 Depósito": "所需租金",
    "800MXN": "S/.250",
    "☑️ Comisión": "佣金比例",
    "9%": "10%",
    "📈 Ingreso Diario": "每日收益",
    "18-20MXN": "S/.5.1-S/.5.5",
    "💰 Ingreso Anual": "年度收益",
    "540-600MXN": "S/.1861.5-S/.2007.5",
    "Características": "特点",
    "Fácil Inicio": "轻松起步",
    "Baja Barrera de Entrada": "低门槛进入",
    "Bajo Riesgo": "低风险",
    "Depósito Mínimo Requerido": "最低保证金要求",
    "Crecimiento Firme": "稳定增长",
    "Retornos Sostenibles": "持续稳定收益",
}),

# ═══ Slide 12: Junior Store → B2 ═══
(12, {
    "Plan de Tienda Junior": "B2店铺收益计划",
    "Tienda Junior (Equilibrio Óptimo)": "均衡投资，更高回报",
    "🔐 Depósito": "所需租金",
    "2400MXN": "S/.800",
    "☑️ Comisión": "佣金比例",
    "11%": "11%",
    "📈 Ingreso Diario": "每日收益",
    "188-205MXN": "S/.17.5-S/.18.2",
    "💰 Ingreso Anual": "年度收益",
    "68620-74825MXN": "S/.6387.5-S/.6643",
    "Ventajas del Plan": "计划优势",
    "✔ Inversión Equilibrada": "✔ 均衡投资",
    "✔ Mayor Visibilidad": "✔ 提高可见度",
    "✔ Mayor Tasa de Comisión": "✔ 更高佣金比例",
    "✔ Soporte Prioritario": "✔ 优先支持",
}),

# ═══ Slide 13: Media Store → I1 ═══
(13, {
    "Plan de Tienda Media": "I1店铺收益计划",
    "Tienda Media": "高级资格，顶级收益",
    "🔐 Depósito Requerido": "所需租金",
    "7000 MXN": "S/.2400",
    "☑️ Comisión": "佣金比例",
    "13%": "13%",
    "💰Ingreso diario: 188-205 / día": "每日收益：S/.55-S/.56",
    "💰 Ingreso anual: 68620-74825 / año": "年度收益：S/.20075-S/.20440",
    "Ventajas del Plan": "计划优势",
    "Estado Premium": "高级资格",
    "Beneficios VIP": "顶级福利",
    "Comisión Max": "最高佣金率",
    "13% de comisión": "13%佣金",
    "Soporte Exclusivo": "专属支持",
    "Equipo dedicado": "敬业团队",
    "Máximo Ingreso": "最高收入",
    "Protección de ingresos": "收益保障",
}),

# ═══ Slide 14: Team Collaboration ═══
(14, {
    "Colaboración en Equipo": "团队协作",
    "Tema: Crear mayor valor personal a través del equipo": "通过团队合作创造更大的个人价值",
    "Concepto Central: En Swihub, nadie trabaja aislado. A través de la colaboración eficiente, cada miembro aporta valor y l": "核心理念\n在Bonroy Apex Sync，没有人是孤立工作的。\n通过高效协作，每一位成员在为平台贡献价值的同时，也在实现个人财富增长与社会价值提升。",
    "Elementos Centrales:": "核心要素",
    "Colaboración": "协作",
    "Trabajo en equipo eficiente": "高效协同工作",
    "Crecimiento de Riqueza": "财富增长",
    "Aumentar ingresos personales": "提升个人收入",
    "Red Global": "全球网络",
    "Conectar recursos globales": "连接全球资源",
    "Valor Social": "社会价值",
    "Crear impacto positivo": "创造积极影响",
}),

# ═══ Slide 15: Expansion → South America ═══
(15, {
    "Expansión América N/S": "南美洲扩张",
    "Estrategia de Expansión: Estableceremos oficinas en ciudades principales de América. Estas oficinas gestionarán ops regi": "扩张战略\n未来将在南美洲主要城市设立办公室\n这些办公室将负责区域运营管理，扩大品牌影响力，并在吸引本地人才的同时创造就业机会。",
    "Funciones Centrales:": "核心职能",
    "Ops Regionales: Gestión y soporte local": "区域运营：本地化管理与支持",
    "Expansión de Marca: Mejorar influencia en el mercado": "品牌扩张：提升市场影响力",
    "Creación de Empleo: Oportunidades de empleo local": "就业创造：为当地创造就业机会",
    "Diseño Regional:": "区域布局",
    "Región Norteamérica": "北美区域",
    "Ciudad de México (Centro LatAm), Ottawa (Centro Estratégico de Canadá), Nueva York (Centro Financiero Global / Sede), Lo": "墨西哥城（México）拉美核心枢纽\n渥太华（Canada）加拿大战略中心\n纽约（USA）全球资本与金融中心\n洛杉矶（USA）西海岸商业与物流中心\n多伦多（Canada）加拿大经济与商业核心",
    "Sudamérica": "南美洲",
    "Bogotá (Colombia), Medellín (Colombia), Manizales (Colombia), Lima (Perú), São Paulo (Brasil)": "波哥大（哥伦比亚）麦德林（哥伦比亚）马尼萨莱斯（哥伦比亚）利马（秘鲁）圣保罗（巴西）",
}),

# ═══ Slide 16: Revenue Overview ═══
(16, {
    "Visión General del Sistema de Ingresos": "收益系统概览",
    "1. Niveles de Tienda e Ingresos": "1. 店铺等级与收益",
    "2. Recompensas por Referidos (Pago Único)": "2. 联盟奖励（一次性）",
    "3. Comisión por Tareas de Equipo (Continuo)": "3. 团队任务佣金（持续收益）",
}),

# ═══ Slide 17: Rewards System ═══
(17, {
    "Sistema de Recompensas": "联盟奖励系统",
    "Nivel A (Directo): \n8% \nReferencia Directa": "A级（直推）\nB1,B2:5%  I1:7%",
    "Nivel B: \n4% \nSegundo Nivel": "B级\nB1,B2:3%  I1:4%",
    "Nivel C: \n2% \nTercer Nivel": "C级\nB1,B2:1%  I1:2%",
    "Ejemplo de Cálculo": "奖励计算示例",
    "Invite a un socio de Tienda Media y reciba: \nRecompensa instantánea de \n\n\n      560 MXN\n\n\n\n\n\nMétodo de cálculo: \n\n7000 *": "邀请一位B2店铺合作伙伴，可获得：\nS/.40即时奖励\n\n计算方式：\n800*5%=S/.40",
}),

# ═══ Slide 18: Commission ═══
(18, {
    "Ingresos por Comisiones": "团队任务佣金收入",
    "Ingresos pasivos de las tareas del equipo: \nCuando su equipo completa tareas de la plataforma, usted recibe una comisión": "来自团队的被动收入\n当您的团队成功完成平台任务时，您将获得佣金\n这将基于团队活跃度，形成可持续的被动收入来源",
    "Estructura de Tasa de Comisión": "佣金结构",
    "Tarea Nivel A\nTareas de equipo directo — 5%\n\nTarea Nivel B\nTareas de segundo nivel — 3%\n\nTarea Nivel C\nTareas de tercer ": "A级\n直接团队任务 — 3%\n\nB级\n第二层级任务 — 2%\n\nC级\n第三层级任务 — 1%",
    "Escenario de Ejemplo": "示例场景",
    "Socio Medio (Nivel A)": "B2合作伙伴（A级）",
    "Ingreso Pasivo Diario": "每日被动收入",
    "Aprox. \n9.4 - 10.25 MXN/dia": "约 S/.0.525-S/.0.546/天",
    # Shape 5 labels (need to check)
}),

# ═══ Slide 19: Priority Strategy ═══
(19, {
    "Estrategia de Prioridad": "排名优先策略",
    "Restricciones de Reglas": "关键规则",
    "Límite de Recompensa\nSi invita a un socio con un nivel superior a su nivel actual, su recompensa por invitación estará l": "奖励限制\n如果加入您合作社的伙伴等级高于您当前等级，\n您的联盟奖励将以您当前店铺等级为上限。",
    "Consejo Estratégico y Sistema de Niveles": "战略建议",
    "Le recomendamos subir de nivel su tienda para maximizar las ganancias del equipo y desbloquear todo el potencial de comi": "我们建议您提升店铺等级，以最大化团队收益，\n并解锁完整的佣金潜力。",
}),

# ═══ Slide 20: Regional Partner ═══
(20, {
    "Socio Regional y Gerente Regional": "区域合作伙伴与区域经理",
    "Nivel 1: Socio Regional": "一级：区域合作伙伴",
    "Requisito:": "要求",
    "15 socios de Nivel A\n(No más de 10 propietarios de Tienda Micro entre los socios de Nivel A)": "合作社有10名A级合作伙伴",
    "Salario (Fijo):": "月薪",
    "3,600 MXN": "S/.200",
    "Nivel 2: Gerente Regional": "二级：区域经理",
    "15 socios de Nivel A\n(El equipo incluye 2 Socios Regionales)": "合作社有25个合作伙伴，其中至少有10个A级合作伙伴",
    "7,200 MXN": "S/.500",
    "Ruta de Ascenso: → Construya su equipo → Gane salario fijo → Suba nivel": "晋升路径 → 建立你的团队 → 赚取固定工资 → 提升等级",
}),

# ═══ Slide 21: City Partner ═══
(21, {
    "Socio de Ciudad y Gerente de Ciudad": "城市合作伙伴与城市主管",
    "Nivel 3: Socio de Ciudad": "三级：城市合作伙伴",
    "El equipo incluye 4 Gerentes Regionales": "合作社包含50个合作伙伴，其中至少有10个A级合作伙伴",
    "27,000 MXN": "S/.1200",
    "Nivel 4: Gerente de Ciudad": "四级：城市主管",
    "30 socios de Nivel A\n(El equipo incluye 3 Socios de Ciudad)": "合作社有100个合作伙伴，其中至少有15个A级合作伙伴",
    "72,000 MXN": "S/.2500",
    # Bottom labels
    "Formación de Equipo": "团队建设",
    "Liderazgo": "领导力",
    "Gestión": "管理能力",
    "Honor": "荣誉认可",
}),

# ═══ Slide 22: City Director ═══
(22, {
    "Promoción | Director de Ciudad y Gerente Prov.": "晋升 | 城市总监与省级代理",
    "Nivel 5: Director de Ciudad": "五级：城市总监",
    "El equipo incluye 2 Gerentes de Ciudad": "合作社有300个合作伙伴，其中至少有15个A级合作伙伴",
    "144,000 MXN": "S/.6000",
    "Nivel 6: Agente Provincial": "六级：省级代理",
    "El equipo incluye 3 Gerentes de Ciudad": "合作社有800个合作伙伴，其中至少有15个A级合作伙伴",
    "271,000 MXN": "S/.16000",
    "Crecimiento de Liderazgo: \nDe Director de Ciudad a Agente Provincial. Salario Máx": "领导力成长路径\n从城市总监到省级代理\n每月最高薪资 — S/.16000",
}),

# ═══ Slide 23: National Agent ═══
(23, {
    "Agente Nacional": "全国代理",
    "Nivel Superior: El Pináculo del Éxito": "顶级等级：成功的巅峰",
    "Salario Fijo Mensual Superior": "顶级月度固定薪资",
    "543,000 MXN": "S/.38000",
    "Requisito: El equipo incluye 3 Agentes Provinciales": "要求：合作社有2000个合作伙伴，其中至少有15个A级合作伙伴",
    "Cobertura Nacional | Gran Red | Estado Élite | Máximos Retornos": "全国覆盖 | 大型网络 | 精英地位 | 最高收益",
}),

# ═══ Slide 24: AI Smart Hosting ═══
(24, {
    "Sistema de Alojamiento Inteligente AI": "AI智能托管系统",
    "Ops Automatizadas 24/7: \n\nLa tecnología libera el trabajo humano y aumenta la eficiencia.": "24/7小时自动化运营\n\n技术不仅解放人力，还通过智能自动化大幅提升效率。",
    "Siempre en Línea": "始终在线",
    "Operaciones 7x24 ininterrumpidas": "7×24小时不间断运行",
    "Velocidad de Milisegundos": "毫秒级速度",
    "Emparejamiento instantáneo": "即时匹配",
    "Cobertura Global": "全球覆盖",
    "Soporta pedidos globales": "支持全球订单",
    "IA Inteligente": "智能AI",
    "Apoyo a la toma de decisiones": "智能决策支持",
}),

# ═══ Slide 25: T+ Settlement ═══
(25, {
    "Motor de Liquidación Ultra Rápido T+": "T+极速结算引擎",
    "Posicionamiento: Liquidación transparente y distribución rápida": "透明结算与快速收益分配",
    "Compromiso Central / Características:": "核心特性",
    "Transparente": "透明",
    "Registros de transacciones claros": "清晰的交易记录",
    "Rápido": "快速",
    "Liquidación instantánea": "即时结算",
    "Seguridad": "安全",
    "Asegura los ingresos del usuario": "保障用户收益安全",
    "Soporte Técnico: \nImpulsado por interfaces financieras de Nueva York, las transacciones pueden lograr una liquidación in": "技术支持\n由纽约金融接口驱动，交易可实现即时结算\n以机构级速度与安全性保障用户收益",
}),

# ═══ Slide 26: Risk Control ═══
(26, {
    "Control de Riesgos Institucional": "机构级风险控制",
    "Mecanismo: Con amplia experiencia en gestión de activos, aseguramos: \nBajo estándares de Seguridad Institucional, lograr": "凭借丰富的资产管理经验，我们确保：\n在机构级安全标准下，实现透明运营与实时财富追踪。",
    "Capacidades Centrales:": "核心能力",
    "Operaciones Transparentes": "透明运营",
    "Todos los procesos completamente visibles": "所有流程全程可视",
    "Seguimiento en Tiempo Real": "实时追踪",
    "Monitoreo de activos 7x24": "7×24小时监控资产变化",
    "Seguridad Institucional": "机构级安全",
    "Protección de seguridad de grado profesional": "专业级安全防护",
}),

# ═══ Slide 27: Sophia Welcome ═══
(27, {
    "Líder del Proyecto: Sophia | Educación: Universidad de Londres": "项目负责人：Sophia | 教育背景：伦敦大学",
    "Discurso de Bienvenida:": "欢迎致辞",
    "“¡Bienvenido a Swihub! Soy Sophia. Con mi formación de EMBA en la Universidad de Londres, soy responsable personalmente ": "欢迎来到 Bonroy Apex Sync！我是 Sophia，凭借我在伦敦大学获得的EMBA背景，我亲自负责为新参与者打造数字资产增长计划。我们的团队将全程陪伴您，在迈向财务增长与数字成功的每一步提供支持与指导。",
    "Soporte Central:\n\nOrientación de Expertos | Soporte Personalizado | Orientado al Crecimiento": "核心支持\n\n专家指导 | 个性化支持 | 成长导向",
}),

# ═══ Slide 28: Internship ═══
(28, {
    "Fase de Dueño de Tienda en Prácticas (Días 1–4)": "实习店主阶段（第1–4天）",
    "Prueba de 6 Días | Sin Depósito Requerido": "4天试用 | 无需保证金",
    "Duración": "周期",
    "4 Días": "4天",
    "Artículos": "商品数量",
    "6 Artículos": "3件商品",
    "Ingreso Diario": "每日收益",
    "17.5-19 MXN/Día": "S/.4.5-S/.5/单日",
    "Plataforma de Experiencia: \nAntes de hacer cualquier compromiso, experimente nuestro flujo de trabajo y sistema de liqui": "体验平台\n在做出任何承诺之前，以零风险体验我们的工作流程与结算系统。",
    "Experimente Ahora": "立即体验",
}),

# ═══ Slide 29: First Withdrawal ═══
(29, {
    "Primer Retiro": "首次提现",
    "Contacte a sus amigos para obtener orientación: Complete su primer retiro de 35 MXN\n\nConstruyendo confianza: Cuando los ": "联系朋友获得指引\n完成您的首次 S/.10 提现\n\n建立信任\n当用户体验到真实收益时，信任便开始建立。\n首次提现是一个重要里程碑，体现了我们兑现承诺的能力。",
    "Proceso de Retiro:": "提现流程",
    "1. Completar Tarea": "1. 完成任务",
    "4 dias Gane al menos 70-76 MXN": "至少赚取S/.18-S/.20",
    "2. Enviar Solicitud": "2. 提交提现申请",
    "Enviar solicitud de retiro": "提交提现请求",
    "3. Recibir Fondos": "3. 收到款项",
    "Fondos llegan al instante": "资金即时到账",
    "Retiro Mínimo: 35 MXN": "最低提现金额：S/.10",
}),

# ═══ Slide 30: FAQ & Support ═══
(30, {
    "Preguntas Frecuentes y Soporte 24/7": "常见问题 & 24/7支持",
    "El equipo de Sophia brinda soporte continuo": "Sophia团队提供全天候支持",
    "Contenido de Soporte": "支持内容",
    "Cumplimiento\n     Cuestiones regulatorias\n\nSeguridad\n     Protección de cuenta\n\nRetiro\n     Soporte de pagos y retiros\n\n": "合规\n     监管相关问题\n\n安全\n     账户保护\n\n提现\n     支付与提现支持\n\n常见问题\n     平台使用指南",
    "Canales de Soporte": "支持渠道",
    "Chat en Vivo:\nSoporte de mensajería instantánea\n\n\nSoporte por Correo:\nConsultas detalladas": "在线聊天\n即时消息支持（24/7）\n\n邮件支持\n详细问题咨询（24/7）",
    "Siempre Aquí para Ayudar": "始终为您提供帮助",
}),


# ═══ Slide 31: Closing ═══
(31, {
    "En Swihub, los participantes no solo obtienen retornos, sino también el orgullo de contribuir a la sociedad.": "在Bonroy Apex Sync，参与者不仅获得财务回报，还将收获为社会做出贡献的自豪感。",
    "Retornos Financieros": "财务回报",
    "Logre un crecimiento de ingresos\nsostenible a través del comercio digital": "通过数字商业实现可持续收入增长",
    "Impacto Social": "社会影响",
    "Hacer contribuciones sustanciales\na la comunidad y causas públicas": "为社区与公益事业做出贡献",
    "Logro Personal": "个人成就感",
    "Sea parte de una causa\nsignificativa a nivel global": "成为全球有意义事业的一部分",
    "“Únase a nosotros para remodelar juntos el futuro del comercio digital. \nCreemos valor, acumulemos riqueza y generemos u": "加入我们，共同重塑数字商业的未来。\n携手创造价值、积累财富，并在全球范围内带来积极的社会影响。",
    "Swihub — Crecimiento de ingresos digitales a su alcance": "Bonroy Apex Sync — 让数字收入增长触手可及",
    "📍 Ubicado en Rockefeller Plaza, Midtown Manhattan, listo para conectar el futuro del comercio digital global—esperamos q": "📍 英国，伦敦——我们期待您的加入。",
}),
]

print(f"Template has {len(prs.slides)} slides")
print(f"Mapping has {len(slides_map)} slides to create")

# Track which slides were created (indices after original 32)
new_indices = []

# Process each mapping
for orig_idx, mapping in slides_map:
    ns = duplicate_slide(prs, orig_idx)
    slide_num = len(prs.slides) - 1
    new_indices.append(slide_num)
    replace_text_on_slide(ns, mapping)
    # Debug: show first text shape
    for s in ns.shapes:
        if s.has_text_frame and s.text_frame.text.strip():
            print(f"  Slide {slide_num} (from {orig_idx}): '{s.text_frame.text.strip()[:60]}'")
            break

print(f"Created {len(new_indices)} new slides. Total: {len(prs.slides)}")

# ── Table fills (only on our new slides) ──
# Slides 16 (49) → 3 tables
s16 = prs.slides[new_indices[16]]
# Table 0: Store levels
fill_table(s16, "Table 0", [
    (0,0,"等级"), (0,1,"租金(索尔)"), (0,2,"商品数量"), (0,3,"佣金"),
    (0,4,"每日收益"), (0,5,"每月收益"), (0,6,"每年收益"),
    (1,0,"B1店铺"), (1,1,"250"), (1,2,"4个"), (1,3,"10%"),
    (1,4,"5.1-5.5"), (1,5,"153-165"), (1,6,"1861.5-2007.5"),
    (2,0,"B2店铺"), (2,1,"800"), (2,2,"6个"), (2,3,"11%"),
    (2,4,"17.5-18.2"), (2,5,"525-546"), (2,6,"6387.5-6643"),
    (3,0,"I1店铺"), (3,1,"2400"), (3,2,"8个"), (3,3,"13%"),
    (3,4,"55-56"), (3,5,"1650-1680"), (3,6,"20075-20440"),
])
# Table 1: Alliance incentives
fill_table(s16, "Table 1", [
    (0,0,"店铺等级"), (0,1,"比例"), (0,2,"A级（直接）"), (0,3,"B级（二层）"), (0,4,"C级（三层）"),
    (1,0,"B1"), (1,1,"(5%-3%-1%)"), (1,2,"S/.12.5"), (1,3,"S/.7.5"), (1,4,"S/.2.5"),
    (2,0,"B2"), (2,1,"(5%-3%-1%)"), (2,2,"S/.40"), (2,3,"S/.24"), (2,4,"S/.8"),
    (3,0,"I1"), (3,1,"(7%-4%-2%)"), (3,2,"S/.168"), (3,3,"S/.96"), (3,4,"S/.48"),
])
# Table 2: Team commission
fill_table(s16, "Table 2", [
    (0,0,"店铺等级"), (0,1,"比例"), (0,2,"A级/天"), (0,3,"B级/天"), (0,4,"C级/天"),
    (1,0,"B1"), (1,1,"(3%-2%-1%)"), (1,2,"0.153-0.165"), (1,3,"0.102-0.11"), (1,4,"0.051-0.055"),
    (2,0,"B2"), (2,1,"(3%-2%-1%)"), (2,2,"0.525-0.546"), (2,3,"0.35-0.364"), (2,4,"0.175-0.182"),
    (3,0,"I1"), (3,1,"(4%-3%-1%)"), (3,2,"2.2-2.24"), (3,3,"1.65-1.68"), (3,4,"0.55-0.56"),
])

# Slide 17 (50): replace specific shapes not caught by text matching
s17 = prs.slides[new_indices[17]]
for s in s17.shapes:
    if s.name in ('文本框 14', '文本框 15', '文本框 16') and s.has_text_frame:
        labels = {'文本框 14': 'A', '文本框 15': 'B', '文本框 16': 'C'}
        if s.name in labels:
            s.text = labels[s.name]

# Slide 18 (51): replace shape labels
s18 = prs.slides[new_indices[18]]
for s in s18.shapes:
    if s.has_text_frame:
        t = s.text_frame.text.strip()
        if t == "Level A\n5%": s.text = "A级\n3%"
        elif t == "Level B\n3%": s.text = "B级\n2%"
        elif t == "Level C\n1%": s.text = "C级\n1%"
    # Rectangles labeled A/B/C
    if s.name in ('矩形 12', '矩形 13', '矩形 14') and s.has_text_frame:
        labels = {'矩形 12': 'A', '矩形 13': 'B', '矩形 14': 'C'}
        if s.name in labels:
            s.text = labels[s.name]

# Slide 19 (52): catch multi Text 7 shapes (card labels)
s19 = prs.slides[new_indices[19]]
for s in s19.shapes:
    if s.has_text_frame:
        t = s.text_frame.text.strip()
        if "Tienda Junior" in t or "12%" in t:
            s.text = "B1店铺：10%佣金"
        elif "Tienda Media" in t or "10%" in t:
            s.text = "B2店铺：11%佣金"
        elif "Tienda Senior" in t or "8%" in t:
            s.text = "I1店铺：13%佣金"

# ── Save temp, then rebuild clean zip (removes orphan slide files) ──
# WARNING: python-pptx can't safely delete slides — direct XML manipulation
# leaves orphan slide files in the zip, breaking Office compatibility.
# Use zip-level cleanup instead.
tmp_path = OUTPUT + ".tmp.pptx"
prs.save(tmp_path)
print(f"Temp saved: {len(prs.slides)} slides")

import zipfile, os as _os
from lxml import etree as _et

with zipfile.ZipFile(tmp_path, 'r') as zin:
    all_files = zin.namelist()
    pres_xml = _et.fromstring(zin.read('ppt/presentation.xml'))
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
    sldIdLst = pres_xml.find(f'.//{{{ns_p}}}sldIdLst')
    all_sld_ids = sldIdLst.findall(f'{{{ns_p}}}sldId')
    # Keep only the last 32 (the new slides), remove first 32 (originals)
    keep_rIds = set()
    remove_rIds = set()
    for i, si in enumerate(all_sld_ids):
        rid = si.get(f'{{{ns_r}}}id')
        if i >= 32:
            keep_rIds.add(rid)
        else:
            remove_rIds.add(rid)
            sldIdLst.remove(si)
    # Map rIds to slide files
    srels = _et.fromstring(zin.read('ppt/_rels/presentation.xml.rels'))
    rels_ns = 'http://schemas.openxmlformats.org/package/2006/relationships'
    keep_parts = set()
    remove_parts = set()
    for rel in srels.findall(f'{{{rels_ns}}}Relationship'):
        rid = rel.get('Id')
        target = rel.get('Target')
        if 'slide' in target.lower():
            sn = 'ppt/slides/' + target.split('/')[-1]
            if rid in keep_rIds:
                keep_parts.add(sn)
                keep_parts.add('ppt/slides/_rels/' + target.split('/')[-1] + '.rels')
            elif rid in remove_rIds:
                remove_parts.add(sn)
                remove_parts.add('ppt/slides/_rels/' + target.split('/')[-1] + '.rels')
    # Write clean zip
    with zipfile.ZipFile(OUTPUT, 'w', zipfile.ZIP_DEFLATED) as zout:
        for f in all_files:
            if f in remove_parts: continue
            if f == '[Content_Types].xml':
                ct = _et.fromstring(zin.read(f))
                ct_ns = 'http://schemas.openxmlformats.org/package/2006/content-types'
                for ov in list(ct.findall(f'{{{ct_ns}}}Override')):
                    pn = ov.get('PartName')
                    if '/slides/slide' in pn and ('ppt/slides/' + pn.rsplit('/',1)[-1]) in remove_parts:
                        ct.remove(ov)
                zout.writestr(f, _et.tostring(ct, xml_declaration=True, encoding='UTF-8', standalone=True))
            elif f == 'ppt/presentation.xml':
                zout.writestr(f, _et.tostring(pres_xml, xml_declaration=True, encoding='UTF-8', standalone=True))
            elif f == 'ppt/_rels/presentation.xml.rels':
                prels = _et.fromstring(zin.read(f))
                for rel in list(prels.findall(f'{{{rels_ns}}}Relationship')):
                    if rel.get('Id') in remove_rIds:
                        prels.remove(rel)
                zout.writestr(f, _et.tostring(prels, xml_declaration=True, encoding='UTF-8', standalone=True))
            else:
                zout.writestr(f, zin.read(f))
_os.remove(tmp_path)
print(f"Clean PPTX: {len(keep_rIds)} slides saved to {OUTPUT}")
print("DONE")
