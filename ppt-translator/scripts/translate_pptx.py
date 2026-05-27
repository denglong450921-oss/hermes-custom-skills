import json
import math
import os
import re
import sys
import time
import subprocess
from collections import Counter
from datetime import datetime
from pathlib import Path

# --- Dedicated Virtual Environment Bootstrap ---
# This ensures the skill runs with its required dependencies without polluting the global environment
# or requiring the user to run `pip install` in every new project.
SCRIPT_DIR = Path(__file__).resolve().parent
VENV_DIR = SCRIPT_DIR.parent / "venv"
VENV_PYTHON = VENV_DIR / "bin" / "python3"

def _ensure_venv():
    if os.environ.get("PPT_TRANSLATOR_VENV_ACTIVE") == "1":
        return

    if sys.executable == str(VENV_PYTHON):
        return

    ready_marker = VENV_DIR / ".ready"
    if not ready_marker.exists():
        print(f"[ppt-translator] Initializing dedicated virtual environment...", flush=True)
        if VENV_DIR.exists():
            import shutil
            shutil.rmtree(VENV_DIR)
        subprocess.run([sys.executable, "-m", "venv", str(VENV_DIR)], check=True)
        
        pip_path = VENV_DIR / "bin" / "pip"
        print("[ppt-translator] Installing dependencies (python-pptx, deep-translator)...", flush=True)
        env = os.environ.copy()
        for k in ["http_proxy", "https_proxy", "all_proxy", "HTTP_PROXY", "HTTPS_PROXY", "ALL_PROXY"]:
            env.pop(k, None)
        subprocess.run(
            [str(pip_path), "install", "--proxy", "", "-i", "https://pypi.tuna.tsinghua.edu.cn/simple", "python-pptx", "deep-translator"],
            check=True,
            env=env
        )
        ready_marker.touch()
        print("[ppt-translator] Environment ready.", flush=True)

    os.environ["PPT_TRANSLATOR_VENV_ACTIVE"] = "1"
    os.environ["PATH"] = f"{VENV_DIR / 'bin'}:{os.environ.get('PATH', '')}"
    os.execv(str(VENV_PYTHON), [str(VENV_PYTHON)] + sys.argv)

_ensure_venv()
# -----------------------------------------------

import pptx
from deep_translator import GoogleTranslator
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.text.fonts import FontFiles
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree

TERMS = [
    "Prism AI Systems Inc",
    "Prism Systems, Inc.",
    "Prism AI",
    "PAS",
    "B2B",
    "B2C",
]

# Patterns where a whole text block should skip translation entirely
# All-caps abbreviations with numbers (B2B, SaaS, v1.0), version strings, model codes
_SKIP_PATTERNS = [
    re.compile(r'^[\d\W]+$'),               # digits only or punctuation only (original)
    re.compile(r'^[A-Z0-9\s\-_./+#(),;:&%]+$'),  # all-caps + digits + symbols (B2B B2C, SaaS, v1.0)
    re.compile(r'^[\d]+[kKmMbBgGtTpP]?[WwHhZz]?$'),   # pure units (100kW, 50Hz, 10TB)
]

# Direct Google Translate API (bypasses deep-translator's outdated lang list)
_GT_URL = "https://translate.googleapis.com/translate_a/single"

def _gt_direct_api(text, target_lang, retries=3):
    """Translate using raw Google Translate API. Supports languages deep-translator doesn't list."""
    import requests as _req, urllib.parse as _up
    for i in range(retries):
        try:
            r = _req.get(
                f"{_GT_URL}?client=gtx&sl=auto&tl={target_lang}&dt=t&q={_up.quote(text)}",
                timeout=10,
                headers={"User-Agent": "Mozilla/5.0"}
            )
            if r.status_code == 200:
                return "".join(part[0] for part in r.json()[0] if part[0])
        except Exception:
            time.sleep(2)
    return text

SCRIPT_DIR = Path(__file__).resolve().parent
CACHE_DIR = Path.home() / ".cache" / "ppt-translator"
CACHE_FILE = CACHE_DIR / "translation_cache.json"
translation_cache = {}
FONT_FALLBACKS = [
    "Arial Unicode MS",
    "Arial",
]
RUNTIME_STATS = {
    "translated_blocks": 0,
    "cache_hits": 0,
    "skipped_blocks": 0,
}
LANGUAGE_ALIASES = {
    "malayalam": "ml",
    "ml": "ml",
    "马拉雅拉姆语": "ml",
    "telugu": "te",
    "te": "te",
    "泰卢固语": "te",
    "tswana": "tn",
    "tn": "tn",
    "茨瓦纳语": "tn",
    "spanish": "es",
    "es": "es",
    "西班牙语": "es",
    "french": "fr",
    "fr": "fr",
    "法语": "fr",
    "german": "de",
    "de": "de",
    "德语": "de",
    "japanese": "ja",
    "ja": "ja",
    "日语": "ja",
    "korean": "ko",
    "ko": "ko",
    "韩语": "ko",
    "arabic": "ar",
    "ar": "ar",
    "阿拉伯语": "ar",
    "hindi": "hi",
    "hi": "hi",
    "印地语": "hi",
    "tamil": "ta",
    "ta": "ta",
    "泰米尔语": "ta",
    "marathi": "mr",
    "mr": "mr",
    "马拉地语": "mr",
    "bengali": "bn",
    "bn": "bn",
    "孟加拉语": "bn",
    "gujarati": "gu",
    "gu": "gu",
    "古吉拉特语": "gu",
    "kannada": "kn",
    "kn": "kn",
    "卡纳达语": "kn",
    "punjabi": "pa",
    "pa": "pa",
    "旁遮普语": "pa",
    "urdu": "ur",
    "ur": "ur",
    "乌尔都语": "ur",
    "portuguese": "pt",
    "pt": "pt",
    "葡萄牙语": "pt",
    "russian": "ru",
    "ru": "ru",
    "俄语": "ru",
    "italian": "it",
    "it": "it",
    "意大利语": "it",
    "indonesian": "id",
    "id": "id",
    "印尼语": "id",
    "thai": "th",
    "th": "th",
    "泰语": "th",
    "vietnamese": "vi",
    "vi": "vi",
    "越南语": "vi",
}

# Right-to-left (RTL) languages — text direction must flip for these
_RTL_LANGUAGES = {"ar", "he", "iw", "fa", "ur", "ps", "sd", "ku", "yi", "dv"}


def log(message):
    """Print a timestamped progress log so long-running translation stays visible."""
    timestamp = datetime.now().strftime("%H:%M:%S")
    print(f"[ppt-translator {timestamp}] {message}", flush=True)

def load_cache():
    """Load the shared translation cache stored alongside the bundled script."""
    global translation_cache
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    if CACHE_FILE.exists():
        with CACHE_FILE.open('r', encoding='utf-8') as f:
            translation_cache = json.load(f)
        log(f"Loaded translation cache: {len(translation_cache)} entries")
    else:
        log("No existing translation cache found, starting fresh")

def save_cache():
    """Persist translated segments so repeated PPT runs stay fast."""
    CACHE_DIR.mkdir(parents=True, exist_ok=True)
    with CACHE_FILE.open('w', encoding='utf-8') as f:
        json.dump(translation_cache, f, ensure_ascii=False, indent=2)
    log(f"Saved translation cache: {len(translation_cache)} entries")

def protect_terms(text):
    """Replace protected terms with placeholders before machine translation."""
    term_map = {}
    for i, term in enumerate(TERMS):
        placeholder = f"[[TERM{i}]]"
        if term in text:
            text = text.replace(term, placeholder)
            term_map[placeholder] = term
    return text, term_map

def restore_terms(text, term_map):
    """Restore protected terms after translation is complete."""
    for placeholder, term in term_map.items():
        text = text.replace(placeholder, term)
    return text

def translate_with_retry(text, target_lang, retries=4):
    """Translate one text block with cache lookup and retry-based backoff."""
    if not text.strip(): 
        RUNTIME_STATS["skipped_blocks"] += 1
        return text
    stripped = text.strip()
    if any(p.match(stripped) for p in _SKIP_PATTERNS): 
        RUNTIME_STATS["skipped_blocks"] += 1
        return text
        
    cache_key = f"{target_lang}::{text}"
    if cache_key in translation_cache:
        RUNTIME_STATS["cache_hits"] += 1
        return translation_cache[cache_key]
    
    protected_text, term_map = protect_terms(text)
    
    for i in range(retries):
        try:
            translated = _gt_direct_api(protected_text, target_lang, retries=1)
            time.sleep(0.5) 
            final_text = restore_terms(translated, term_map)
            translation_cache[cache_key] = final_text
            RUNTIME_STATS["translated_blocks"] += 1
            return final_text
        except Exception as e:
            log(f"Translation retry {i + 1}/{retries} for one text block: {e}")
            time.sleep(2.0) 
            
    log(f"WARNING: failed to translate one text block after {retries} retries")
    return restore_terms(protected_text, term_map)

def get_text_frames(shapes):
    """Yield text frames recursively from normal shapes, groups, and tables."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from get_text_frames(shape.shapes)
        elif shape.has_text_frame:
            yield shape.text_frame
        elif shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        yield cell.text_frame


def count_translatable_blocks(presentation):
    """Count non-empty paragraph blocks so the user can see expected workload."""
    count = 0
    for slide in presentation.slides:
        for text_frame in get_text_frames(slide.shapes):
            for paragraph in text_frame.paragraphs:
                full_text = "".join(run.text for run in paragraph.runs).strip()
                if full_text:
                    count += 1
    return count


def get_run_font_size_pt(run, default_size=18.0):
    """Return the run font size in points, falling back to a safe default."""
    if run.font.size is not None:
        return run.font.size.pt
    return default_size


def get_run_font_name(run):
    """Return the run font family when explicitly defined."""
    return run.font.name


def resolve_measurement_font(paragraph_entries):
    """Pick a measurable font for fitting that is as close to the source as possible."""
    family_counter = Counter(
        entry["font_name"] for entry in paragraph_entries if entry["font_name"]
    )
    families = [family for family, _ in family_counter.most_common()]
    families.extend(FONT_FALLBACKS)

    for family in families:
        try:
            return family, FontFiles.find(family, False, False)
        except Exception:
            continue

    return "Arial", FontFiles.find("Arial", False, False)


def apply_frame_fit(text_frame, paragraph_entries):
    """Compute one frame-level scale factor and preserve each paragraph's relative size."""
    if not paragraph_entries:
        return

    max_original_size = max(entry["orig_size_pt"] for entry in paragraph_entries)
    if max_original_size <= 0:
        return

    family, font_file = resolve_measurement_font(paragraph_entries)

    try:
        best_fit_size = text_frame._best_fit_font_size(
            family=family,
            max_size=max(1, math.floor(max_original_size)),
            bold=False,
            italic=False,
            font_file=font_file,
        )
    except Exception:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        return

    scale_factor = min(1.0, best_fit_size / max_original_size)
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.word_wrap = True

    for entry in paragraph_entries:
        scaled_size = max(1.0, entry["orig_size_pt"] * scale_factor)
        entry["first_run"].font.size = Pt(scaled_size)


def apply_rtl_layout(text_frame):
    """Set right-to-left text direction for the entire text frame and all paragraphs."""
    # Body-level RTL column
    bodyPr = text_frame._bodyPr
    bodyPr.set('rtlCol', '1')

    # Each paragraph: RTL + right alignment
    a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    for para in text_frame.paragraphs:
        para.alignment = PP_ALIGN.RIGHT
        pPr = para._p.find(qn('a:pPr'))
        if pPr is None:
            pPr = etree.SubElement(para._p, qn('a:pPr'))
        pPr.set('rtl', '1')

        # Also set rtl on endParaRPr if present
        endRPr = para._p.find(qn('a:endParaRPr'))
        if endRPr is not None:
            rPr_child = endRPr.find(qn('a:rPr'))
            if rPr_child is None:
                rPr_child = etree.SubElement(endRPr, qn('a:rPr'))
            rPr_child.set('rtl', '1')


def is_rtl_target(lang_code):
    """Check if the normalized language code needs RTL layout."""
    return lang_code in _RTL_LANGUAGES


def normalize_target_language(target_lang):
    """Normalize human-readable language names to translator-friendly codes."""
    normalized = LANGUAGE_ALIASES.get(target_lang.strip().lower())
    if normalized:
        return normalized
    return target_lang.strip().lower()


def build_output_path(input_path, target_lang):
    """Generate a predictable translated PPTX filename when one is not provided."""
    source_path = Path(input_path)
    normalized_lang = normalize_target_language(target_lang)
    safe_suffix = re.sub(r"[^A-Za-z0-9_-]+", "_", normalized_lang).strip("_") or "translated"
    return source_path.with_name(f"{source_path.stem}_{safe_suffix}{source_path.suffix}")

def process_presentation(input_path, output_path, target_lang):
    """Translate an entire presentation while preserving layout and text hierarchy."""
    for key in RUNTIME_STATS:
        RUNTIME_STATS[key] = 0

    normalized_target = normalize_target_language(target_lang)
    log(f"Starting translation: source='{input_path}', target='{normalized_target}'")
    is_rtl = is_rtl_target(normalized_target)
    if is_rtl:
        log("RTL language detected — applying right-to-left text direction")
    load_cache()
    prs = pptx.Presentation(input_path)
    slide_total = len(prs.slides)
    total_blocks = count_translatable_blocks(prs)
    log(f"Presentation loaded: {slide_total} slides, {total_blocks} text blocks to inspect")
    
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_blocks = 0
        for text_frame in get_text_frames(slide.shapes):
            text_frame.word_wrap = True
            paragraph_entries = []

            for paragraph in text_frame.paragraphs:
                runs = paragraph.runs
                if not runs:
                    continue

                full_text = "".join(r.text for r in runs)
                if not full_text.strip():
                    continue

                first_run = runs[0]
                orig_size_pt = None
                for r in runs:
                    if r.font.size is not None:
                        orig_size_pt = r.font.size.pt
                        break
                if orig_size_pt is None:
                    orig_size_pt = get_run_font_size_pt(first_run)

                translated_text = translate_with_retry(full_text, normalized_target)
                first_run.text = translated_text
                for r in runs[1:]:
                    r.text = ""
                paragraph.line_spacing = 1.0
                slide_blocks += 1

                paragraph_entries.append(
                    {
                        "first_run": first_run,
                        "orig_size_pt": orig_size_pt,
                        "font_name": get_run_font_name(first_run),
                    }
                )

            apply_frame_fit(text_frame, paragraph_entries)
            if is_rtl and paragraph_entries:
                apply_rtl_layout(text_frame)
        log(f"Processed slide {slide_index}/{slide_total}: {slide_blocks} text blocks" + (" [RTL]" if is_rtl else ""))

    save_cache()
    prs.save(output_path)
    log(
        "Finished translation: "
        f"translated={RUNTIME_STATS['translated_blocks']}, "
        f"cache_hits={RUNTIME_STATS['cache_hits']}, "
        f"skipped={RUNTIME_STATS['skipped_blocks']}"
    )
    log(f"Saved output presentation: {output_path}")
    return output_path

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python translate_pptx.py <input.pptx> <target_lang> [output.pptx]")
        sys.exit(1)
        
    input_file = sys.argv[1]
    target = sys.argv[2]
    output = sys.argv[3] if len(sys.argv) >= 4 else str(build_output_path(input_file, target))
    final_output = process_presentation(input_file, output, target)
    print(f"OUTPUT_FILE={final_output}")
